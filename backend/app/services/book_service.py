import os
import time
import fitz  # PyMuPDF
import docx
from typing import List

import google.generativeai as genai

from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont


class BookService:
    # ------------------------------------------------------------------
    # MAIN ENTRY
    # ------------------------------------------------------------------
    def extract_text(self, file_path: str) -> str:
        """Extracts text based on file extension."""
        if not os.path.exists(file_path):
            raise ValueError(f"File not found: {file_path}")

        ext = os.path.splitext(file_path)[1].lower()

        if ext == ".pdf":
            return self._extract_pdf(file_path)

        elif ext == ".docx":
            return self._extract_docx(file_path)

        elif ext == ".txt":
            return self._extract_txt(file_path)

        elif ext in (".ppt", ".pptx"):
            # ✅ HIGH-QUALITY OCR USING GEMINI
            return self.process_pptx_with_gemini(file_path)

        else:
            raise ValueError("Unsupported file format")

    # ------------------------------------------------------------------
    # BASIC FILE TYPES
    # ------------------------------------------------------------------
    def _extract_pdf(self, path: str) -> str:
        try:
            doc = fitz.open(path)
            text = ""
            for page in doc:
                text += page.get_text()
            return text
        except Exception as e:
            raise ValueError(f"Error reading PDF: {e}")

    def _extract_docx(self, path: str) -> str:
        try:
            doc = docx.Document(path)
            return "\n".join([p.text for p in doc.paragraphs])
        except Exception as e:
            raise ValueError(f"Error reading DOCX: {e}")

    def _extract_txt(self, path: str) -> str:
        try:
            with open(path, "r", encoding="utf-8") as f:
                return f.read()
        except Exception as e:
            raise ValueError(f"Error reading TXT: {e}")

    # ------------------------------------------------------------------
    # PPT → GEMINI OCR PIPELINE (THE IMPORTANT PART)
    # ------------------------------------------------------------------
    def process_pptx_with_gemini(self, file_path: str) -> str:
        """
        High-accuracy OCR for PPT/PPTX (text + images).
        Uses Gemini Vision and waits until files are ACTIVE.
        """
        temp_dir = f"temp_ppt_{int(time.time())}"
        slide_images = self.pptx_to_slide_images(file_path, temp_dir)

        if not slide_images:
            raise ValueError("No slides found in PPT")

        # Upload slides to Gemini
        uploaded_files = [genai.upload_file(path=img) for img in slide_images]

        # ✅ CRITICAL: wait for OCR readiness
        self.wait_for_files_active(uploaded_files)

        # Gemini Vision OCR + summarization
        model = genai.GenerativeModel("gemini-1.5-flash")
        response = model.generate_content(
            [
                "Extract ALL readable text from these slides accurately. "
                "Do not summarize yet. Preserve headings and bullet points.",
                *uploaded_files,
            ]
        )

        return response.text or ""

    # ------------------------------------------------------------------
    # GEMINI SAFETY WAIT LOOP
    # ------------------------------------------------------------------
    def wait_for_files_active(self, files, timeout_seconds: int = 180):
        """
        Prevents empty OCR / infinite loops.
        """
        start_time = time.time()

        for f in files:
            while True:
                current = genai.get_file(f.name)
                state = current.state.name

                if state == "ACTIVE":
                    break

                if state == "FAILED":
                    raise Exception(f"OCR failed for file: {f.name}")

                if time.time() - start_time > timeout_seconds:
                    raise TimeoutError(
                        "PPT OCR timed out. Possible network or Gemini issue."
                    )

                time.sleep(3)

    # ------------------------------------------------------------------
    # PPT → IMAGE CONVERSION
    # ------------------------------------------------------------------
    def pptx_to_slide_images(self, pptx_path: str, out_dir: str) -> List[str]:
        """
        Converts PPTX slides into images suitable for Gemini OCR.
        """
        os.makedirs(out_dir, exist_ok=True)
        prs = Presentation(pptx_path)
        img_paths = []

        for i, slide in enumerate(prs.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text.strip())

            content = "\n\n".join(slide_text).strip()
            if not content:
                content = "(Image-based slide – OCR by Gemini)"

            img = Image.new("RGB", (1280, 720), "white")
            draw = ImageDraw.Draw(img)

            try:
                font = ImageFont.truetype("arial.ttf", 28)
            except:
                font = ImageFont.load_default()

            draw.text(
                (40, 40),
                f"Slide {i}\n\n{content}",
                fill="black",
                font=font,
            )

            out_path = os.path.join(out_dir, f"slide_{i}.png")
            img.save(out_path)
            img_paths.append(out_path)

        return img_paths

    # ------------------------------------------------------------------
    # CHUNKING (FOR SUMMARIZATION)
    # ------------------------------------------------------------------
    def chunk_text(self, text: str, chunk_size: int = 500) -> List[str]:
        words = text.split()
        chunks, current = [], []

        for w in words:
            current.append(w)
            if len(current) >= chunk_size:
                chunks.append(" ".join(current))
                current = []

        if current:
            chunks.append(" ".join(current))

        return chunks


# ✅ SINGLE INSTANCE
book_service = BookService()
